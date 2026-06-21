"""
课件智析 CourseMind - Flask 后端主程序
======================================
功能：
  1. 接收文件上传（PDF / PPT）
  2. 提供前端静态页面
  3. 后续将接入 DeepSeek AI
"""

import os
import sys
import json
from datetime import datetime
from pathlib import Path

# === PDF 解析库 ===
from PyPDF2 import PdfReader

# === PPT/PPTX 解析库 ===
from pptx import Presentation

# === 修复 Windows 中文/Emoji 编码问题 ===
if sys.platform == "win32":
    try:
        import io
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")
    except Exception:
        pass  # 有些终端不支持，忽略即可

from flask import Flask, request, jsonify, send_from_directory

# ==========================================
# 1. 初始化 Flask 应用
# ==========================================

# 当前文件所在目录就是 backend/
BASE_DIR = Path(__file__).parent

# 告诉 Flask：前端文件在 ../frontend/ 目录
FRONTEND_DIR = BASE_DIR.parent / "frontend"

# 上传文件保存目录
UPLOAD_DIR = BASE_DIR / "uploads"
UPLOAD_DIR.mkdir(exist_ok=True)  # 确保目录存在

# 允许上传的文件类型
ALLOWED_EXTENSIONS = {"pdf", "ppt", "pptx"}

app = Flask(__name__, static_folder=str(FRONTEND_DIR), static_url_path="")

# Flask 配置
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024  # 最大上传 100MB


# ==========================================
# 2. 辅助函数
# ==========================================

# --- DeepSeek AI 配置 ---
# 可通过环境变量覆盖：DEEPSEEK_BASE_URL, DEEPSEEK_API_KEY, DEEPSEEK_MODEL
DEEPSEEK_BASE_URL = os.environ.get("DEEPSEEK_BASE_URL", "https://api.deepseek.com")
DEEPSEEK_API_KEY = os.environ.get("DEEPSEEK_API_KEY", "")  # 必须通过环境变量设置
DEEPSEEK_MODEL  = os.environ.get("DEEPSEEK_MODEL", "deepseek-v4-flash")

# 已解析的课件文本缓存（用于 AI 问答上下文）
_parsed_cache: dict = {}  # {filename: full_text}

def allowed_file(filename: str) -> bool:
    """检查文件后缀是否允许"""
    if "." not in filename:
        return False
    ext = filename.rsplit(".", 1)[1].lower()
    return ext in ALLOWED_EXTENSIONS


def get_file_info(filepath: Path) -> dict:
    """获取文件信息"""
    stat = filepath.stat()
    return {
        "name": filepath.name,
        "size": stat.st_size,
        "size_mb": round(stat.st_size / (1024 * 1024), 2),
        "upload_time": datetime.fromtimestamp(stat.st_ctime).strftime("%Y-%m-%d %H:%M:%S"),
    }


def call_deepseek(prompt: str, system_prompt: str = None, temperature: float = 0.7) -> str:
    """
    调用 DeepSeek API（兼容 OpenAI 格式）
    参数：
        prompt: 用户提示词
        system_prompt: 系统提示词（可选）
        temperature: 温度参数（0-1）
    返回：AI 生成的文本
    """
    import requests

    messages = []
    if system_prompt:
        messages.append({"role": "system", "content": system_prompt})
    messages.append({"role": "user", "content": prompt})

    headers = {
        "Content-Type": "application/json",
    }
    if DEEPSEEK_API_KEY and DEEPSEEK_API_KEY != "ollama":
        headers["Authorization"] = f"Bearer {DEEPSEEK_API_KEY}"

    payload = {
        "model": DEEPSEEK_MODEL,
        "messages": messages,
        "temperature": temperature,
        "max_tokens": 4096,
    }

    resp = requests.post(
        f"{DEEPSEEK_BASE_URL}/chat/completions",
        headers=headers,
        json=payload,
        timeout=120,
    )
    resp.raise_for_status()
    result = resp.json()
    return result["choices"][0]["message"]["content"]


# ==========================================
# 3. 路由：前端页面
# ==========================================

@app.route("/")
def index():
    """返回主页面"""
    return send_from_directory(str(FRONTEND_DIR), "index.html")


@app.route("/css/<path:filename>")
def serve_css(filename):
    """返回 CSS 文件"""
    return send_from_directory(str(FRONTEND_DIR / "css"), filename)


@app.route("/js/<path:filename>")
def serve_js(filename):
    """返回 JS 文件"""
    return send_from_directory(str(FRONTEND_DIR / "js"), filename)


# ==========================================
# 4. 核心路由：文件上传
# ==========================================

@app.route("/api/upload", methods=["POST"])
def upload_file():
    """
    接收前端上传的文件
    POST /api/upload
    参数：files（一个或多个文件）
    返回：上传成功的文件列表
    """
    # 检查是否有文件
    if "files" not in request.files:
        return jsonify({"error": "没有收到文件"}), 400

    files = request.files.getlist("files")
    if not files or files[0].filename == "":
        return jsonify({"error": "文件为空"}), 400

    uploaded_files = []
    errors = []

    for file in files:
        # 检查文件名是否合法
        if file.filename == "":
            continue

        # 检查文件类型
        if not allowed_file(file.filename):
            errors.append(f"{file.filename}：不支持的文件类型")
            continue

        try:
            # 安全文件名（防止路径穿越攻击）
            safe_name = Path(file.filename).name
            save_path = UPLOAD_DIR / safe_name

            # 如果同名文件已存在，加时间戳
            if save_path.exists():
                stem = save_path.stem
                ext = save_path.suffix
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                save_path = UPLOAD_DIR / f"{stem}_{timestamp}{ext}"

            # 保存文件
            file.save(str(save_path))

            # 记录文件信息
            info = get_file_info(save_path)
            info["saved_name"] = save_path.name
            uploaded_files.append(info)

            print(f"[上传成功] {save_path.name} ({info['size_mb']} MB)")

        except Exception as e:
            errors.append(f"{file.filename}：保存失败 - {str(e)}")
            print(f"[上传失败] {file.filename}: {e}")

    # 返回结果
    result = {
        "success": len(uploaded_files) > 0,
        "uploaded": uploaded_files,
        "errors": errors,
        "message": f"成功上传 {len(uploaded_files)} 个文件"
        if uploaded_files
        else "上传失败",
    }

    if not uploaded_files:
        return jsonify(result), 400

    return jsonify(result), 200


# ==========================================
# 5. 路由：查看已上传文件列表
# ==========================================

@app.route("/api/files", methods=["GET"])
def list_files():
    """返回已上传的文件列表"""
    files = []
    for f in UPLOAD_DIR.iterdir():
        if f.is_file():
            files.append(get_file_info(f))
    # 按上传时间倒序排列
    files.sort(key=lambda x: x["upload_time"], reverse=True)
    return jsonify({"files": files})


# ==========================================
# 6. 路由：删除已上传文件
# ==========================================

@app.route("/api/files/<filename>", methods=["DELETE"])
def delete_file(filename):
    """删除指定文件"""
    file_path = UPLOAD_DIR / filename
    if not file_path.exists():
        return jsonify({"error": "文件不存在"}), 404
    try:
        file_path.unlink()
        print(f"[删除成功] {filename}")
        return jsonify({"message": f"已删除 {filename}"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


# ==========================================
# 7. PDF 文本提取
# ==========================================

def extract_text_from_pdf(filepath: Path) -> str:
    """
    从 PDF 文件中提取所有文字
    参数：PDF 文件路径
    返回：提取的文本内容
    """
    text_parts = []
    try:
        reader = PdfReader(str(filepath))
        total_pages = len(reader.pages)
        print(f"[PDF解析] 共 {total_pages} 页")

        for i, page in enumerate(reader.pages):
            page_text = page.extract_text()
            if page_text:
                # 清理多余空白
                cleaned = page_text.strip().replace("\n\n", "\n")
                text_parts.append(f"--- 第 {i+1} 页 ---\n{cleaned}")

        full_text = "\n\n".join(text_parts)
        print(f"[PDF解析] 提取文字 {len(full_text)} 个字符")
        return full_text

    except Exception as e:
        print(f"[PDF解析错误] {e}")
        raise


def extract_text_from_pptx(filepath: Path) -> str:
    """
    从 PPT/PPTX 文件中提取所有文字
    参数：PPT 文件路径
    返回：提取的文本内容
    """
    text_parts = []
    try:
        prs = Presentation(str(filepath))
        total_slides = len(prs.slides)
        print(f"[PPT解析] 共 {total_slides} 页")

        for i, slide in enumerate(prs.slides):
            slide_texts = []

            for shape in slide.shapes:
                # 提取文本框中的文字
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            slide_texts.append(t)

                # 提取表格中的文字
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_texts = []
                        for cell in row.cells:
                            ct = cell.text.strip()
                            if ct:
                                row_texts.append(ct)
                        if row_texts:
                            slide_texts.append(" | ".join(row_texts))

            if slide_texts:
                text_parts.append(f"--- 第 {i+1} 页 ---\n" + "\n".join(slide_texts))

        full_text = "\n\n".join(text_parts)
        print(f"[PPT解析] 提取文字 {len(full_text)} 个字符")
        return full_text

    except Exception as e:
        print(f"[PPT解析错误] {e}")
        raise


# ==========================================
# 8. 路由：解析已上传的课件
# ==========================================

@app.route("/api/parse", methods=["POST"])
def parse_courseware():
    """
    解析指定的课件文件
    POST /api/parse
    请求体 JSON：{"files": ["文件名1.pdf", "文件名2.pptx"]}
    返回：每个文件的解析结果
    """
    data = request.get_json(silent=True) or {}
    requested_files = data.get("files", [])

    # 如果没有指定文件，回退到解析全部文件（兼容旧版前端）
    if not requested_files:
        requested_files = [f.name for f in UPLOAD_DIR.iterdir() if f.is_file()]

    results = []
    errors = []

    for filename in requested_files:
        filepath = UPLOAD_DIR / filename

        if not filepath.exists():
            errors.append(f"{filename}: 文件不存在")
            continue

        ext = filepath.suffix.lower()

        try:
            if ext == ".pdf":
                text = extract_text_from_pdf(filepath)
                _parsed_cache[filename] = text  # 缓存供 AI 使用
                results.append({
                    "file": filename,
                    "type": "pdf",
                    "pages": len(PdfReader(str(filepath)).pages),
                    "text": text,
                    "text_length": len(text),
                })
                print(f"[解析成功] PDF: {filename}")

            elif ext in [".ppt", ".pptx"]:
                text = extract_text_from_pptx(filepath)
                _parsed_cache[filename] = text  # 缓存供 AI 使用
                prs = Presentation(str(filepath))
                results.append({
                    "file": filename,
                    "type": "ppt",
                    "pages": len(prs.slides),
                    "text": text,
                    "text_length": len(text),
                })
                print(f"[解析成功] PPT: {filename}")

            else:
                errors.append(f"{filename}: 不支持的文件类型")

        except Exception as e:
            errors.append(f"{filename}: 解析失败 - {str(e)}")
            print(f"[解析失败] {filename}: {e}")

    if not results:
        return jsonify({
            "success": False,
            "message": "没有可解析的文件，请先上传课件",
            "results": [],
            "errors": errors,
        }), 400

    return jsonify({
        "success": True,
        "message": f"成功解析 {len(results)} 个文件",
        "results": results,
        "errors": errors,
    })


# ==========================================
# 9. AI 分析路由（DeepSeek）
# ==========================================

def _get_course_text(filename: str = None) -> str:
    """获取课件文本：优先用缓存，否则取第一个缓存文件"""
    if filename and filename in _parsed_cache:
        return _parsed_cache[filename]
    if _parsed_cache:
        return list(_parsed_cache.values())[0]
    return ""


def _build_system_prompt() -> str:
    return (
        "你是一位经验丰富的大学教师，擅长将复杂的课程内容转化为"
        "清晰、结构化的学习资料。请用中文回答，格式整洁有条理。"
    )


@app.route("/api/chat", methods=["POST"])
def chat():
    """
    AI 问答（基于已解析的课件内容）
    POST /api/chat
    请求体 JSON：{"question": "问题", "filename": "可选-指定课件"}
    返回：{"answer": "AI回答"}
    """
    data = request.get_json(silent=True) or {}
    question = data.get("question", "").strip()
    filename = data.get("filename", None)

    if not question:
        return jsonify({"error": "问题不能为空"}), 400

    course_text = _get_course_text(filename)
    if not course_text:
        return jsonify({
            "answer": "⚠️ 请先上传并解析课件，我才能基于课件内容回答你的问题。"
        }), 200

    # 截断过长的文本（保留前 8000 字）
    context = course_text[:8000]

    prompt = f"""请根据以下课件内容回答学生的问题。

【课件内容】
{context}

【学生问题】
{question}

请基于课件内容作答。如果课件中没有相关信息，请如实告知，不要编造。"""

    try:
        answer = call_deepseek(prompt, _build_system_prompt(), temperature=0.5)
        return jsonify({"answer": answer})
    except Exception as e:
        print(f"[AI问答错误] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/ai/summary", methods=["POST"])
def ai_summary():
    """
    AI 课程总结（概述 + 核心知识点 + 章节关系）
    POST /api/ai/summary
    请求体 JSON：{"filename": "可选-指定课件"}
    返回：{"overview": "概述", "key_points": [...], "relations": "..."}
    """
    data = request.get_json(silent=True) or {}
    filename = data.get("filename")

    course_text = _get_course_text(filename)
    if not course_text:
        return jsonify({"error": "请先解析课件"}), 400

    context = course_text[:6000]

    system = _build_system_prompt()
    prompt = f"""请对以下课件内容做课程总结，按以下三个部分输出（用 Markdown 格式）：

## 📖 课程概述
用 2-3 句话概括这门课的主要内容。

## 🔑 核心知识点
列出 5-8 个核心知识点，每个用编号列出，包含简要解释。

## 🔗 章节关系
分析各章节/知识点之间的逻辑关系（递进、并列、因果等）。

【课件内容】
{context}"""

    try:
        answer = call_deepseek(prompt, system, temperature=0.5)
        # 拆分三个部分
        parts = answer.split("## ")
        overview = ""
        key_points = ""
        relations = ""
        for part in parts:
            part_stripped = part.strip()
            if part_stripped.startswith("📖") or part_stripped.startswith("课程概述"):
                overview = part_stripped
            elif part_stripped.startswith("🔑") or part_stripped.startswith("核心知识点"):
                key_points = part_stripped
            elif part_stripped.startswith("🔗") or part_stripped.startswith("章节关系"):
                relations = part_stripped

        return jsonify({
            "overview": overview or "（未能生成概述）",
            "key_points": key_points or "（未能提取知识点）",
            "relations": relations or "（未能分析关系）",
            "raw": answer,
        })
    except Exception as e:
        print(f"[AI总结错误] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/ai/exam", methods=["POST"])
def ai_exam():
    """
    AI 考试重点提炼（★★★★★ 评级）
    POST /api/ai/exam
    返回：{"points": [{"text": "...", "stars": 5, "reason": "..."}, ...]}
    """
    data = request.get_json(silent=True) or {}
    filename = data.get("filename")

    course_text = _get_course_text(filename)
    if not course_text:
        return jsonify({"error": "请先解析课件"}), 400

    context = course_text[:6000]

    prompt = f"""请从以下课件内容中提炼考试重点，按五星评级标注重要性。

对每个考点给出：
- 考点名称
- 星级（★★★★★=必考，★★★★☆=重点，★★★☆☆=了解）
- 简短理由

用以下格式输出（每个考点用 --- 分隔）：
考点：xxx
星级：★★★★★
理由：xxx
---

【课件内容】
{context}"""

    try:
        answer = call_deepseek(prompt, _build_system_prompt(), temperature=0.4)
        return jsonify({"content": answer})
    except Exception as e:
        print(f"[AI考试重点错误] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/ai/framework", methods=["POST"])
def ai_framework():
    """
    AI 知识框架（树状结构）
    POST /api/ai/framework
    返回：{"content": "Markdown 嵌套列表"}
    """
    data = request.get_json(silent=True) or {}
    filename = data.get("filename")

    course_text = _get_course_text(filename)
    if not course_text:
        return jsonify({"error": "请先解析课件"}), 400

    context = course_text[:6000]

    prompt = f"""请将以下课件内容整理成树状知识框架，用 Markdown 嵌套列表表示。

要求：
- 用缩进列表表示层级关系
- 第一层是大章节/主题
- 第二层是知识点
- 第三层是关键细节

【课件内容】
{context}"""

    try:
        answer = call_deepseek(prompt, _build_system_prompt(), temperature=0.4)
        return jsonify({"content": answer})
    except Exception as e:
        print(f"[AI知识框架错误] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/ai/mindmap", methods=["POST"])
def ai_mindmap():
    """
    AI 思维导图（Markdown 格式）
    POST /api/ai/mindmap
    返回：{"content": "Markdown 思维导图"}
    """
    data = request.get_json(silent=True) or {}
    filename = data.get("filename")

    course_text = _get_course_text(filename)
    if not course_text:
        return jsonify({"error": "请先解析课件"}), 400

    context = course_text[:6000]

    prompt = f"""请将以下课件内容整理成思维导图格式，用 Markdown 标题 + 列表表示。

要求：
- # 课程名（中心主题）
- ## 各大模块
- - 知识点（用无序列表）
- 结构清晰、逻辑分明

【课件内容】
{context}"""

    try:
        answer = call_deepseek(prompt, _build_system_prompt(), temperature=0.5)
        return jsonify({"content": answer})
    except Exception as e:
        print(f"[AI思维导图错误] {e}")
        return jsonify({"error": str(e)}), 500


@app.route("/api/ai/quiz", methods=["POST"])
def ai_quiz():
    """
    AI 练习题生成
    POST /api/ai/quiz
    返回：{"content": "选择题 + 判断题 + 简答题，含答案"}
    """
    data = request.get_json(silent=True) or {}
    filename = data.get("filename")

    course_text = _get_course_text(filename)
    if not course_text:
        return jsonify({"error": "请先解析课件"}), 400

    context = course_text[:6000]

    prompt = f"""请根据以下课件内容生成练习题，包含：

## 选择题（5道，每题4个选项，标注正确答案）
## 判断题（3道，标注对/错及理由）
## 简答题（2道，附参考答案要点）

【课件内容】
{context}"""

    try:
        answer = call_deepseek(prompt, _build_system_prompt(), temperature=0.6)
        return jsonify({"content": answer})
    except Exception as e:
        print(f"[AI练习题错误] {e}")
        return jsonify({"error": str(e)}), 500


# ==========================================
# 10. 启动服务器
# ==========================================

if __name__ == "__main__":
    print("=" * 50)
    print("🧠 课件智析 CourseMind 后端服务")
    print("=" * 50)
    print(f"📂 上传目录：{UPLOAD_DIR}")
    print(f"🌐 访问地址：http://localhost:5000")
    print(f"📤 上传接口：POST http://localhost:5000/api/upload")
    print("=" * 50)

    # 生产环境下请设置 FLASK_ENV=production，使用 gunicorn 启动
    is_debug = os.environ.get("FLASK_ENV", "development") == "development"
    app.run(host="0.0.0.0", port=int(os.environ.get("PORT", 5000)), debug=is_debug)
